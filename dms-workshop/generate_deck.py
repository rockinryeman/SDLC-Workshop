#!/usr/bin/env python3
"""Generate the branded AI-Powered Product Development (DMS) slide deck.
Run with: /tmp/pptxvenv/bin/python generate_deck.py
Brand: Industrial DevOps Now — teal #256E8E, dark #174C63, light #3288AD, gold #F0A500."""

import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

# ---------- brand ----------
TEAL  = RGBColor(0x25, 0x6E, 0x8E)
TEALD = RGBColor(0x17, 0x4C, 0x63)
TEALL = RGBColor(0x32, 0x88, 0xAD)
GOLD  = RGBColor(0xF0, 0xA5, 0x00)
INK   = RGBColor(0x1B, 0x1B, 0x1B)
GRAY  = RGBColor(0x59, 0x59, 0x59)
LIGHT = RGBColor(0xF7, 0xF7, 0xF7)
LINE  = RGBColor(0xE2, 0xE2, 0xE2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUB   = RGBColor(0xD8, 0xE6, 0xF7)
FONT  = "Arial"

tTEALD = (0x17, 0x4C, 0x63)
tTEALL = (0x32, 0x88, 0xAD)
def lerp(a, b, t):
    return RGBColor(int(a[0]+(b[0]-a[0])*t), int(a[1]+(b[1]-a[1])*t), int(a[2]+(b[2]-a[2])*t))
def shades(n):
    return [lerp(tTEALD, tTEALL, i/(max(n-1,1))) for i in range(n)]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5

def slide():
    return prs.slides.add_slide(BLANK)

def shape(s, kind, x, y, w, h, fill=None, line=None, line_w=None):
    sp = s.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w or 1)
    return sp

def settext(sp, text, size, color, bold=False, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.MIDDLE, font=FONT, italic=False):
    tf = sp.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        f = r.font
        f.size = Pt(size); f.bold = bold; f.italic = italic
        f.name = font; f.color.rgb = color
    return sp

def textbox(s, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP, font=FONT, italic=False):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    settext(tb, text, size, color, bold, align, anchor, font, italic)
    return tb

def bg(s, color):
    shape(s, MSO_SHAPE.RECTANGLE, -0.1, -0.1, SW + 0.2, SH + 0.2, fill=color)

def footer(s, n):
    shape(s, MSO_SHAPE.RECTANGLE, 0.7, 7.04, 11.93, 0.015, fill=LINE)
    textbox(s, 0.7, 7.08, 7, 0.3, "Industrial DevOps Now  ·  AI-Powered Product Development", 9, GRAY)
    textbox(s, 4.8, 7.08, 4, 0.3, "Mentimeter: menti.com", 9, GRAY, align=PP_ALIGN.CENTER)
    textbox(s, 11.0, 7.08, 1.63, 0.3, str(n), 9, GRAY, align=PP_ALIGN.RIGHT, bold=True)

def accent(s):
    shape(s, MSO_SHAPE.RECTANGLE, 0, 0, 0.16, 7.5, fill=TEAL)

def header(s, title, kicker=None):
    accent(s)
    y = 0.6
    if kicker:
        textbox(s, 0.7, 0.48, 11.5, 0.3, kicker.upper(), 12, GOLD, bold=True)
        y = 0.9
    textbox(s, 0.7, y, 11.9, 0.95, title.upper(), 28, TEALD, bold=True)
    shape(s, MSO_SHAPE.RECTANGLE, 0.72, y + 0.86, 1.1, 0.045, fill=GOLD)
    return y + 1.18

def run_pill(s, text):
    p = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 9.9, 0.62, 2.73, 0.55, fill=GOLD)
    settext(p, "▶  " + text, 13, INK, bold=True, align=PP_ALIGN.CENTER)

def flow(s, items, x, y, total_w, h, fill_list=None, tcolor=WHITE, fsize=12):
    n = len(items); gap = 0.28
    bw = (total_w - gap * (n - 1)) / n
    cur = x
    for i, it in enumerate(items):
        col = fill_list[i] if fill_list else TEAL
        b = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, cur, y, bw, h, fill=col)
        settext(b, it, fsize, tcolor, bold=True, align=PP_ALIGN.CENTER)
        if i < n - 1:
            shape(s, MSO_SHAPE.RIGHT_ARROW, cur + bw + 0.03, y + h/2 - 0.11, gap - 0.06, 0.22, fill=GOLD)
        cur += bw + gap

def hexagon(s, x, y, w, h, fill, label=None, tcolor=WHITE, fsize=12):
    hx = shape(s, MSO_SHAPE.HEXAGON, x, y, w, h, fill=fill)
    if label:
        settext(hx, label, fsize, tcolor, bold=True, align=PP_ALIGN.CENTER)
    return hx

def hex_flow(s, items, x, y, total_w, h, fill_list=None, fsize=12):
    n = len(items); gap = 0.30
    hw = (total_w - gap * (n - 1)) / n
    cur = x
    for i, it in enumerate(items):
        col = fill_list[i] if fill_list else TEAL
        hexagon(s, cur, y, hw, h, col, it, WHITE, fsize)
        if i < n - 1:
            shape(s, MSO_SHAPE.RIGHT_ARROW, cur + hw + 0.01, y + h/2 - 0.10, gap - 0.04, 0.20, fill=GOLD)
        cur += hw + gap

def hex_ring(s, cx, cy, center_label, ring_labels, hw=2.05, hh=1.7, rx=2.75, ry=1.75):
    cols = shades(len(ring_labels))
    for i, lab in enumerate(ring_labels):
        ang = math.radians(-90 + (360/len(ring_labels)) * i)
        x = cx + rx * math.cos(ang) - hw/2
        y = cy + ry * math.sin(ang) - hh/2
        hexagon(s, x, y, hw, hh, cols[i], lab, WHITE, 11)
    hexagon(s, cx - hw/2, cy - hh/2, hw, hh, GOLD, center_label, INK, 13)

def corner_circles(s):
    shape(s, MSO_SHAPE.OVAL, 11.1, -2.3, 5.6, 5.6, fill=TEAL)
    shape(s, MSO_SHAPE.OVAL, 12.1, 3.9, 4.6, 4.6, fill=TEALL)

def set_notes(s, text):
    if text:
        s.notes_slide.notes_text_frame.text = text

# ---------- content slide ----------
def content(n, title, body, kicker=None, pill=None, notes=None):
    s = slide()
    cy = header(s, title, kicker)
    if pill:
        run_pill(s, pill)
    textbox(s, 0.7, cy, 11.6, 1.6, body, 20, INK)
    footer(s, n)
    set_notes(s, notes)
    return s

# ---------- divider slide ----------
def divider(n, kicker, title, subtitle, badge=None):
    s = slide()
    bg(s, TEALD)
    corner_circles(s)
    textbox(s, 0.8, 0.6, 9, 0.4, "INDUSTRIAL DEVOPS NOW", 14, WHITE, bold=True)
    if badge:
        b = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.8, 1.7, 3.7, 0.55, fill=GOLD)
        settext(b, badge, 14, INK, bold=True, align=PP_ALIGN.CENTER)
    if kicker:
        textbox(s, 0.82, 2.55, 9, 0.35, kicker.upper(), 13, GOLD, bold=True)
    shape(s, MSO_SHAPE.RECTANGLE, 0.85, 3.0, 1.6, 0.06, fill=GOLD)
    textbox(s, 0.8, 3.2, 10.0, 1.8, title, 40, WHITE, bold=True)
    textbox(s, 0.8, 4.9, 9.3, 1.4, subtitle, 20, SUB)
    textbox(s, 0.8, 7.05, 2, 0.3, str(n), 10, SUB, bold=True)
    return s

# ---------- menti slide ----------
def menti(n, pollno, kind, question, notes=None):
    s = slide()
    bg(s, LIGHT)
    accent(s)
    b = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.7, 0.75, 4.0, 0.6, fill=GOLD)
    settext(b, "▶  MENTIMETER  ·  POLL " + str(pollno), 14, INK, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, 0.72, 1.65, 11, 0.4, kind.upper(), 13, GRAY, bold=True)
    textbox(s, 0.7, 2.4, 8.7, 2.6, question, 30, TEAL, bold=True)
    qr = shape(s, MSO_SHAPE.RECTANGLE, 10.1, 2.5, 2.3, 2.3, fill=WHITE, line=TEAL, line_w=1.5)
    settext(qr, "QR\njoin code", 13, GRAY, align=PP_ALIGN.CENTER)
    textbox(s, 0.7, 5.3, 9, 0.6, "Join at menti.com — enter the code shown on screen.", 16, INK)
    footer(s, n)
    set_notes(s, notes)
    return s

# ---------- prompt card / prompt slide (visible prompt for follow-along) ----------
def prompt_card(s, x, y, w, h, label, text, body_size=15):
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=LIGHT, line=TEAL, line_w=1.25)
    textbox(s, x + 0.28, y + 0.14, w - 0.55, 0.4, label, 13, TEAL, bold=True)
    shape(s, MSO_SHAPE.RECTANGLE, x + 0.28, y + 0.56, w - 0.56, 0.018, fill=LINE)
    textbox(s, x + 0.28, y + 0.68, w - 0.56, h - 0.85, text, body_size, INK)

def prompt_slide(n, step_title, label, text, ask, notes=None):
    s = slide()
    cy = header(s, step_title, "Act 3 · Live DMS exercise")
    run_pill(s, label.split("—")[0].strip().replace("Prompt", "Run Prompt"))
    prompt_card(s, 0.7, cy, 11.9, 4.5, "▶  " + label + "  — paste into the AI", text, body_size=16)
    footer(s, n)
    set_notes(s, notes)   # notes (P1–P7) already include the "ASK THE ROOM" direction
    return s

# display prompts (concise, audience-facing)
DP1 = ("You are a senior automotive systems engineer. Design a Driver Monitoring System (DMS) that "
       "detects when a driver is drowsy, distracted, or in a medical emergency — and responds appropriately.\n\n"
       "Generate:\n•  Stakeholder needs\n•  Functional requirements\n"
       "•  Non-functional requirements (latency, accuracy, privacy)\n•  Constraints & assumptions\n"
       "•  Acceptance criteria\n\nFormat the output as a requirements specification.")
DP2 = ("Using the approved requirements, create a concise architecture description:\n\n"
       "•  System context diagram (in text)\n•  Major components: sensors, perception/AI, decision logic, "
       "HMI, vehicle interface\n•  Interfaces & data flows\n•  Failure modes\n•  Safety considerations\n"
       "•  Privacy & cybersecurity considerations")
DP3 = ("Using the architecture, generate a simple SysML v2 textual model:\n\n"
       "•  Requirements\n•  Actions / functions (DetectDrowsiness, AssessConfidence, TriggerResponse)\n"
       "•  Interfaces / flows\n•  Major system elements\n\nFocus on readability, not tool-specific correctness.")
DP4 = ("Create:\n\n•  Verification strategy\n•  Functional test cases\n"
       "•  Edge cases (sunglasses, darkness, head turned, passenger interference)\n"
       "•  Negative cases (false-positive drowsiness)\n•  Requirements-to-test traceability matrix\n\n"
       "Include objective pass/fail criteria.")
DP5 = ("Generate:\n\n•  High-level implementation approach\n•  Pseudocode for the detect-assess-respond loop\n"
       "•  Logic flow & data processing steps\n•  Assumptions\n\nExplain how it satisfies the requirements.")
DP6 = ("Review the implementation. Identify:\n\n•  Requirements satisfied\n•  Tests that would pass\n"
       "•  Potential defects\n•  Risks (false alarms, missed detections)\n•  Recommended improvements\n\n"
       "Provide a pass/fail assessment.")
DP7 = ("Create:\n\n•  OTA deployment plan\n•  Rollback strategy\n"
       "•  Monitoring metrics (accuracy, false-alarm rate, driver trust)\n•  Release notes\n"
       "•  Customer impact assessment\n•  Success criteria after deployment")

# =================== PROMPTS (speaker notes) ===================
P1 = ("PROMPT 1 — Requirements  (paste into the AI live)\n\n"
      "You are a senior automotive systems engineer. We are designing a Driver Monitoring System "
      "(DMS) that detects when a driver is drowsy, distracted, or experiencing a medical emergency, "
      "and responds appropriately. Generate:\n- Stakeholder needs\n- Functional requirements\n"
      "- Non-functional requirements (latency, accuracy, privacy)\n- Constraints\n- Assumptions\n"
      "- Acceptance criteria\nFormat the output as a requirements specification.\n\n"
      "ASK THE ROOM: What's missing? What would you challenge?")
P2 = ("PROMPT 2 — Architecture & Design\n\n"
      "Using the approved requirements, create a concise architecture description including:\n"
      "- System context diagram (described in text)\n- Major components (sensors, perception/AI, "
      "decision logic, HMI, vehicle interface)\n- Interfaces and data flows\n- Failure modes\n"
      "- Safety considerations\n- Privacy and cybersecurity considerations\n\n"
      "ASK THE ROOM: Where could this go wrong, and how would the system stay safe?")
P3 = ("PROMPT 3 — SysML v2 Text Model\n\n"
      "Using the architecture, generate a simple SysML v2 textual model including:\n- Requirements\n"
      "- Actions/functions (e.g., DetectDrowsiness, AssessConfidence, TriggerResponse)\n"
      "- Interfaces/flows\n- Major system elements\nFocus on readability, not tool-specific correctness.\n\n"
      "ASK THE ROOM: This is just a structured way to write the design so a person AND a machine can "
      "read it — note how it traces back to the requirements.")
P4 = ("PROMPT 4 — Test Development\n\n"
      "Create:\n- Verification strategy\n- Functional test cases\n- Edge case test cases (sunglasses, "
      "darkness, head turned, passenger interference)\n- Negative test cases (false-positive drowsiness)\n"
      "- A requirements-to-test traceability matrix\nInclude objective pass/fail criteria.\n\n"
      "ASK THE ROOM: Which of these tests would you most want to see pass before shipping?")
P56 = ("PROMPT 5 — Implementation Approach\n\n"
       "Generate:\n- High-level implementation approach\n- Pseudocode for the detect-assess-respond loop\n"
       "- Logic flow\n- Data processing steps\n- Assumptions\nExplain how the solution satisfies the "
       "requirements.\n\nASK: Notice we went from idea to logic with no separate dev-team handoff.\n\n"
       "----------\n\nPROMPT 6 — Verification Review\n\n"
       "Review the proposed implementation. Identify:\n- Which requirements are satisfied\n"
       "- Which tests would pass\n- Potential defects\n- Risks (false alarms and missed detections)\n"
       "- Recommended improvements\nProvide a pass/fail assessment.\n\n"
       "ASK THE ROOM: Where does human expertise still have to make the final call?")
P7 = ("PROMPT 7 — Deployment Planning\n\n"
      "Create:\n- OTA deployment plan\n- Rollback strategy\n- Monitoring metrics (detection accuracy, "
      "false-alarm rate, driver trust signals)\n- Release notes\n- Customer impact assessment\n"
      "- Success criteria after deployment\n\n"
      "ASK THE ROOM: We just reached a deployment plan. How many teams would this normally take?")
PRIME = ("Before Prompt 1, prime the chat once:\n\"We're running a live workshop. Keep each answer "
         "concise and skimmable — bullets over prose.\"")

# =================== BUILD ===================

# 1 — Title (Leidos-style: white, bold uppercase)
s = slide()
textbox(s, 0.85, 0.55, 9, 0.4, "INDUSTRIAL DEVOPS NOW", 14, TEAL, bold=True)
textbox(s, 0.82, 2.6, 11.7, 2.2, "AI-POWERED PRODUCT DEVELOPMENT", 44, TEALD, bold=True)
shape(s, MSO_SHAPE.RECTANGLE, 0.85, 4.8, 1.7, 0.06, fill=GOLD)
textbox(s, 0.85, 5.05, 10.6, 1.0,
        "Building Better Automotive Products with Smaller, Faster, Cross-Functional Teams", 20, GRAY)
textbox(s, 0.85, 6.2, 10, 0.4, "A 90-minute hands-on workshop", 14, TEAL, bold=True)
shape(s, MSO_SHAPE.RECTANGLE, 0.7, 7.02, 11.93, 0.02, fill=GOLD)
textbox(s, 0.7, 7.08, 6, 0.3, "Industrial DevOps Now", 9, GRAY, bold=True)

# 2 — Why we're here
content(2, "Why we're here",
        "A tired driver drifts toward the lane line. Could the car notice — and help?\n\n"
        "Today, building that one feature takes many teams and many months.",
        kicker="Act 1 · Why product development is slow")

# 3 — Objectives (chips)
s = slide()
cy = header(s, "Workshop objectives", "Act 1 · Why product development is slow")
chips = ["Smaller teams", "Faster iteration", "Better traceability", "Outcome-focused delivery"]
cx, cyy, cw, ch, gap = 0.7, cy + 0.3, 5.7, 1.3, 0.5
for i, c in enumerate(chips):
    x = cx + (i % 2) * (cw + gap)
    y = cyy + (i // 2) * (ch + gap)
    b = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch, fill=LIGHT, line=LINE, line_w=1)
    settext(b, c, 22, TEAL, bold=True, align=PP_ALIGN.CENTER)
footer(s, 3)

# 4 — Poll 1
menti(4, 1, "Multiple choice",
      "Where does your organization spend the most engineering effort?")

# 5 — Automotive development today
content(5, "Automotive development today",
        "Specialized functions create handoffs, delays, and rework.",
        kicker="Act 1 · Why product development is slow")

# 6 — The cost of handoffs (chain)
s = slide()
cy = header(s, "The cost of handoffs", "Act 1 · Why product development is slow")
textbox(s, 0.7, cy, 11.6, 0.6, "Every handoff between specialized functions adds delay and a chance for rework.", 18, INK)
hex_flow(s, ["Requirements", "Systems", "Architecture", "Development", "Test", "Release"],
         0.7, 3.5, 11.93, 1.5, fill_list=shades(6), fsize=12)
footer(s, 6)

# 7 — Poll 2
menti(7, 2, "Word cloud", "What slows product delivery the most today?")

# 8 — Why specialization exists
content(8, "Why specialization exists",
        "Risk reduction, compliance, and deep expertise — good reasons that no longer scale alone.",
        kicker="Act 2 · What AI changes")

# 9 — AI as capability amplification
content(9, "AI as capability amplification",
        "AI makes expertise available at the point of need — it augments specialists, not replaces them.",
        kicker="Act 2 · What AI changes")

# 10 — Future-state team (nodes)
s = slide()
cy = header(s, "Future-state team", "Act 2 · What AI changes")
textbox(s, 0.7, cy, 11.6, 0.6, "A small, cross-functional team does what used to need six.", 18, INK)
nodes = [("Product Owner", TEAL, WHITE), ("Engineers", TEAL, WHITE), ("AI Copilot", GOLD, INK)]
nx, ny, nw, nh = 1.1, 3.7, 3.2, 1.3
for i, (label, fill, tc) in enumerate(nodes):
    x = nx + i * (nw + 0.85)
    b = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, ny, nw, nh, fill=fill)
    settext(b, label, 20, tc, bold=True, align=PP_ALIGN.CENTER)
    if i < len(nodes) - 1:
        textbox(s, x + nw, ny, 0.85, nh, "+", 30, TEAL, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
textbox(s, 0.7, 5.4, 11.6, 0.5, "One small, cross-functional product team.", 16, GRAY, italic=True)
footer(s, 10)

# 11 — Scenario divider
divider(11, "Act 3 · Live DMS exercise", "Driver Monitoring System",
        "Detect drowsy, distracted, or medical-emergency states — and respond safely.",
        badge="THE LIVE EXERCISE")

# 12 — Customer need
content(12, "Customer need",
        "Protect a tired or distracted driver — without constant false alarms.",
        kicker="Act 3 · Live DMS exercise")

# 13 — Step 1 Requirements (prompt on slide)
prompt_slide(13, "Step 1 — Requirements", "Prompt 1 — Requirements", DP1,
             "What's missing? What would you challenge?", notes=PRIME + "\n\n" + P1)

# 14 — Step 2 Architecture (prompt on slide)
prompt_slide(14, "Step 2 — Architecture", "Prompt 2 — Architecture", DP2,
             "Where could this go wrong, and how does it stay safe?", notes=P2)

# 15 — Step 3 Model
prompt_slide(15, "Step 3 — Model (SysML v2)", "Prompt 3 — SysML v2 Model", DP3,
             "See how it traces back to the requirements.", notes=P3)

# 16 — Step 4 Tests
prompt_slide(16, "Step 4 — Tests & traceability", "Prompt 4 — Test Development", DP4,
             "Which test would you most want to pass before shipping?", notes=P4)

# 17 — Poll 3
menti(17, 3, "Multiple choice", "Which SDLC activity benefited most from AI so far?")

# 18 — Step 5-6
s = slide()
cy = header(s, "Step 5–6 — Implementation & verification", "Act 3 · Live DMS exercise")
run_pill(s, "Run Prompts 5–6")
prompt_card(s, 0.7, cy, 5.82, 3.85, "▶  Prompt 5 — Implementation", DP5, 13)
prompt_card(s, 6.78, cy, 5.82, 3.85, "▶  Prompt 6 — Verification", DP6, 13)
textbox(s, 0.7, cy + 3.98, 11.9, 0.5, "Ask the room:   Where does human expertise make the final call?",
        14, TEAL, bold=True, italic=True)
footer(s, 18)
set_notes(s, P56)

# 19 — Step 7 Deployment
prompt_slide(19, "Step 7 — Deployment", "Prompt 7 — Deployment Planning", DP7,
             "How many teams would this normally take?", notes=P7)

# 20 — Digital thread (flow)
s = slide()
cy = header(s, "The digital thread", "Act 4 · What just happened")
textbox(s, 0.7, cy, 11.6, 0.6, "One connected, auditable trail — created by a single team in minutes.", 18, INK)
hex_flow(s, ["Need", "Requirement", "Design", "Test", "Code", "Deploy"],
         0.7, 3.5, 11.93, 1.5, fill_list=shades(6), fsize=13)
footer(s, 20)

# 21 — Timeline comparison (bars)
s = slide()
header(s, "Traditional vs. AI-enabled timeline", "Act 4 · What just happened")
# traditional
textbox(s, 0.7, 2.7, 2.6, 0.5, "Traditional", 16, TEAL, bold=True, anchor=MSO_ANCHOR.MIDDLE)
b1 = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 3.3, 2.65, 9.0, 0.85, fill=TEAL)
settext(b1, "Months — sequential handoffs across teams", 15, WHITE, bold=True, align=PP_ALIGN.CENTER)
# ai-enabled
textbox(s, 0.7, 4.3, 2.6, 0.5, "AI-Enabled", 16, GOLD, bold=True, anchor=MSO_ANCHOR.MIDDLE)
b2 = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 3.3, 4.25, 2.6, 0.85, fill=GOLD)
settext(b2, "Hours", 15, INK, bold=True, align=PP_ALIGN.CENTER)
textbox(s, 6.05, 4.3, 6, 0.7, "— one small team", 15, INK, anchor=MSO_ANCHOR.MIDDLE)
textbox(s, 0.7, 5.6, 11.6, 0.6, "Same rigor. A fraction of the calendar time.", 16, GRAY, italic=True)
footer(s, 21)

# 22 — Humans in the loop (honeycomb)
s = slide()
header(s, "Where humans stayed in the loop", "Act 4 · What just happened")
hex_ring(s, 6.65, 4.35, "Humans\nin the loop",
         ["Confidence\nthresholds", "Safety\naggressiveness", "Bias\nchecks",
          "Accountability", "Governance", "Reviews"])
textbox(s, 0.7, 6.5, 11.9, 0.4, "AI assisted. People decided.", 16, GRAY, italic=True, align=PP_ALIGN.CENTER)
footer(s, 22)

# 23 — Poll 4
menti(23, 4, "Ranking", "Where would you pilot AI first?")

# 24 — Future operating model
content(24, "Future operating model",
        "Outcome-oriented product teams, with specialists providing governance, standards, "
        "reviews, and exception handling.",
        kicker="Act 5 · Future-state organization")

# 25 — Key takeaways (cards)
s = slide()
cy = header(s, "Key takeaways", "Act 5 · Future-state organization")
cards = ["Smaller teams", "Faster learning", "Better outcomes"]
cw, ch, gap = 3.7, 1.6, 0.45
cx = 0.7
for i, c in enumerate(cards):
    x = cx + i * (cw + gap)
    b = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, cy + 0.3, cw, ch, fill=shades(3)[i])
    settext(b, c, 22, WHITE, bold=True, align=PP_ALIGN.CENTER)
textbox(s, 0.7, cy + 2.3, 11.6, 0.6, "Humans remain accountable throughout.", 18, INK, bold=True)
footer(s, 25)

# 26 — Closing + Poll 5 (Leidos-style icon divider)
s = slide()
bg(s, TEALD)
shape(s, MSO_SHAPE.OVAL, 5.42, 1.3, 2.5, 2.5, fill=None, line=WHITE, line_w=4)
textbox(s, 5.42, 1.4, 2.5, 2.35, "?", 96, WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
textbox(s, 0, 4.2, 13.333, 0.8, "Questions  ·  Poll 5", 30, WHITE, bold=True, align=PP_ALIGN.CENTER)
textbox(s, 0, 5.15, 13.333, 0.6, "What is the biggest opportunity?   —   join at menti.com", 16, SUB,
        align=PP_ALIGN.CENTER)
textbox(s, 0.7, 7.06, 4, 0.3, "Industrial DevOps Now", 9, SUB, bold=True)
textbox(s, 11.0, 7.06, 1.63, 0.3, "26", 9, SUB, bold=True, align=PP_ALIGN.RIGHT)

out = "AI-Powered-Product-Development_DMS.pptx"
prs.save(out)
print("Saved", out, "with", len(prs.slides._sldIdLst), "slides")
